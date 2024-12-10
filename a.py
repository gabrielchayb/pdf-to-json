import json
from pptx import Presentation

def pptx_to_json(pptx_file, json_file):
    presentation = Presentation(pptx_file)
    slides_data = []

    for slide in presentation.slides:
        slide_data = {"shapes": []}
        for shape in slide.shapes:
            shape_data = {
                "text": shape.text if hasattr(shape, "text") else "",
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }
            slide_data["shapes"].append(shape_data)
        slides_data.append(slide_data)

    with open(json_file, 'w', encoding='utf-8') as f:
        json.dump(slides_data, f, ensure_ascii=False, indent=4)

if __name__ == "__main__":
    pptx_file = "cienciavendas.pptx"
    json_file = "output.json"
    pptx_to_json(pptx_file, json_file)